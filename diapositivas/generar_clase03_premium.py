#!/usr/bin/env python3
"""Generador de diapositivas PREMIUM - Clase 03: Lectura del Precio y Zonas Clave"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
import os

# === COLORES PREMIUM ===
NAVY_DARK = RGBColor(0x0B, 0x1A, 0x2E)
NAVY_MID = RGBColor(0x12, 0x2B, 0x45)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
GOLD_LIGHT = RGBColor(0xF0, 0xC8, 0x56)
TEAL = RGBColor(0x00, 0xC9, 0xA7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xB0, 0xB8, 0xC4)
GRAY_MID = RGBColor(0x6B, 0x7B, 0x8D)
RED_CANDLE = RGBColor(0xE8, 0x4C, 0x3D)
GREEN_CANDLE = RGBColor(0x2E, 0xCC, 0x71)
CARD_BG = RGBColor(0x15, 0x2D, 0x4A)
CARD_BORDER = RGBColor(0x1E, 0x3A, 0x5F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# === HELPER FUNCTIONS ===
def set_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=WHITE, spacing=1.2, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(font_size * spacing * 0.5)
    return txBox


def add_card(slide, left, top, width, height, title, body_lines,
             accent_color=GOLD, title_size=16, body_size=14):
    """Add a card with accent border"""
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = accent_color
    shape.line.width = Pt(1.5)
    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2),
                 width - Inches(0.6), Inches(0.5),
                 title, font_size=title_size, bold=True, color=accent_color)
    # Body
    if body_lines:
        add_multi_text(slide, left + Inches(0.3), top + Inches(0.65),
                       width - Inches(0.6), height - Inches(0.85),
                       body_lines, font_size=body_size, color=GRAY_LIGHT)

def add_candle_strip(slide, y_pos, count=8):
    """Draw decorative candlestick strip"""
    start_x = Inches(8.5)
    spacing = Inches(0.5)
    colors = [GREEN_CANDLE, RED_CANDLE, GREEN_CANDLE, GREEN_CANDLE,
              RED_CANDLE, GREEN_CANDLE, RED_CANDLE, GREEN_CANDLE]
    heights = [Inches(0.6), Inches(0.45), Inches(0.7), Inches(0.5),
               Inches(0.55), Inches(0.8), Inches(0.4), Inches(0.65)]
    for i in range(min(count, len(colors))):
        x = start_x + spacing * i
        h = heights[i]
        # Body
        body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_pos, Inches(0.15), h)
        body.fill.solid()
        body.fill.fore_color.rgb = colors[i]
        body.line.fill.background()
        # Wick
        wick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.06), y_pos - Inches(0.15), Inches(0.03), h + Inches(0.3))
        wick.fill.solid()
        wick.fill.fore_color.rgb = colors[i]
        wick.line.fill.background()


def add_footer(slide, slide_num, total_slides):
    """Add slide number and course name footer"""
    # Slide number
    add_text_box(slide, Inches(0.5), SLIDE_H - Inches(0.5), Inches(1.5), Inches(0.4),
                 f"{slide_num:02d} / {total_slides:02d}", font_size=10, color=GRAY_MID)
    # Course name
    add_text_box(slide, SLIDE_W - Inches(4), SLIDE_H - Inches(0.5), Inches(3.5), Inches(0.4),
                 "Curso de Trading · Índices Sintéticos y Forex", font_size=10,
                 color=GRAY_MID, alignment=PP_ALIGN.RIGHT)

def add_divider_slide(title_text, part_num):
    """Create a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(slide, NAVY_DARK)
    # Big number
    add_text_box(slide, Inches(1), Inches(1.5), Inches(3), Inches(3),
                 f"{part_num:02d}", font_size=120, bold=True, color=CARD_BORDER, font_name='Arial')
    # Gold bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.2), Inches(3), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    # Title
    add_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(1.2),
                 title_text, font_size=32, bold=True, color=WHITE)
    # Candles decoration
    add_candle_strip(slide, Inches(2.5), count=6)
    return slide

TOTAL_SLIDES = 17
slide_counter = [0]

def next_slide():
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(slide, NAVY_DARK)
    add_footer(slide, slide_counter[0], TOTAL_SLIDES)
    return slide


# ============================================================
# SLIDE 1: PORTADA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_solid_bg(slide, NAVY_DARK)
slide_counter[0] += 1

# Gold accent bar left
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), SLIDE_H)
bar.fill.solid()
bar.fill.fore_color.rgb = GOLD
bar.line.fill.background()

# Module label
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(6), Inches(0.5),
             "MÓDULO 1 · ANÁLISIS TÉCNICO", font_size=14, bold=True, color=GOLD)

# Main title
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(8), Inches(1.5),
             "Lectura del Precio\ny Zonas Clave", font_size=48, bold=True, color=WHITE)

# Subtitle
add_text_box(slide, Inches(1.2), Inches(4.5), Inches(6), Inches(0.6),
             "Clase 3 · Índices Sintéticos y Forex", font_size=18, color=TEAL)

# Footer
add_text_box(slide, Inches(1.2), Inches(6.2), Inches(6), Inches(0.5),
             "Curso de Trading desde cero", font_size=12, color=GRAY_MID)

# Candle decoration
add_candle_strip(slide, Inches(2.8), count=8)

# ============================================================
# SLIDE 2: DIVISOR - PARTE 01
# ============================================================
add_divider_slide("Lectura del Precio", 1)
slide_counter[0] += 1


# ============================================================
# SLIDE 3: EL PRECIO COMO LENGUAJE
# ============================================================
slide = next_slide()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "El precio habla un lenguaje", font_size=32, bold=True, color=WHITE)

# Quote box
quote_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.4), Inches(11), Inches(1.5))
quote_shape.fill.solid()
quote_shape.fill.fore_color.rgb = CARD_BG
quote_shape.line.color.rgb = GOLD
quote_shape.line.width = Pt(1)

add_text_box(slide, Inches(1.2), Inches(1.6), Inches(10.2), Inches(1.3),
             '"Leer el precio es como leer un lenguaje: las velas nos cuentan una historia.\n'
             'Cada movimiento, cada rechazo o rompimiento tiene una intención del mercado."',
             font_size=16, color=GOLD_LIGHT, alignment=PP_ALIGN.LEFT)

# Key points
add_multi_text(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(3.5),
    ["▪  HH (Higher High) + HL (Higher Low) = Tendencia ALCISTA",
     "▪  LH (Lower High) + LL (Lower Low) = Tendencia BAJISTA",
     "",
     "▪  El precio NO se mueve al azar → respeta zonas clave",
     "▪  Lo importante no es adivinar, sino entender la secuencia"],
    font_size=18, color=WHITE)

# Speaker notes
slide.notes_slide.notes_text_frame.text = (
    "Leer el precio es como leer un lenguaje: las velas nos cuentan una historia. "
    "Cada movimiento, cada rechazo o rompimiento tiene una intención del mercado.\n\n"
    "Muestra un gráfico (Volatility 75) y di: Vean cómo el precio va dejando huellas. "
    "Cada impulso y cada retroceso forman parte de una estructura.\n\n"
    "Énfasis: El precio no se mueve al azar, respeta zonas clave. "
    "Lo importante no es adivinar, sino entender la secuencia del precio."
)


# ============================================================
# SLIDE 4: DIVISOR - PARTE 02
# ============================================================
add_divider_slide("Soportes y Resistencias", 2)
slide_counter[0] += 1

# ============================================================
# SLIDE 5: QUE SON SOPORTES Y RESISTENCIAS
# ============================================================
slide = next_slide()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Soportes y Resistencias", font_size=32, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.8),
             "Zonas donde el precio ha reaccionado varias veces.\nNo opero al primer toque.",
             font_size=18, color=GRAY_LIGHT)

# Quote
quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(2.3), Inches(11), Inches(1.0))
quote.fill.solid()
quote.fill.fore_color.rgb = CARD_BG
quote.line.color.rgb = TEAL
quote.line.width = Pt(1)

add_text_box(slide, Inches(1.2), Inches(2.45), Inches(10.2), Inches(0.8),
             '"Espero testeos de confirmación. El último es el operativo."',
             font_size=17, bold=True, color=TEAL)

# Notes
slide.notes_slide.notes_text_frame.text = (
    "Los soportes y resistencias son zonas donde el precio ha reaccionado varias veces. "
    "Yo no opero al primer toque. Espero testeos de confirmación; el último es el operativo.\n\n"
    "Concepto clave: No se entra al primer toque. Se espera validación."
)

# ============================================================
# SLIDE 6: LOS 3 TOQUES (TARJETAS)
# ============================================================
slide = next_slide()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Método de los Testeos", font_size=32, bold=True, color=WHITE)

# Three cards for each touch
add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(4.5),
         "1er TOQUE", ["El precio 'descubre'", "la zona.", "", "Primera reacción.", "Aún no operamos."],
         accent_color=GRAY_MID, title_size=18, body_size=15)

add_card(slide, Inches(4.7), Inches(1.5), Inches(3.8), Inches(4.5),
         "2do TOQUE", ["Confirma que el área", "tiene fuerza.", "", "Valida la zona.", "Aún no operamos."],
         accent_color=GOLD, title_size=18, body_size=15)

add_card(slide, Inches(8.9), Inches(1.5), Inches(3.8), Inches(4.5),
         "3er TOQUE ✓", ["OPERATIVO", "", "Busco entrada si", "coincide con:", "• Estructura", "• Hora", "• Confirmación"],
         accent_color=TEAL, title_size=18, body_size=15)

slide.notes_slide.notes_text_frame.text = (
    "Desglosa:\n"
    "• 1er toque: el precio 'descubre' la zona.\n"
    "• 2do toque: confirma que el área tiene fuerza y valida.\n"
    "• 3er toque: ahí busco mi entrada si coincide con estructura, hora y confirmación.\n\n"
    "Importante: cada herramienta (S/R, líneas de tendencia) puede tener su propio criterio de toques."
)


# ============================================================
# SLIDE 7: DIVISOR - PARTE 03
# ============================================================
add_divider_slide("Confirmación con Velas", 3)
slide_counter[0] += 1

# ============================================================
# SLIDE 8: PATRONES + ZONA + HORA
# ============================================================
slide = next_slide()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Patrones de Vela como Confirmación", font_size=32, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "Los patrones NO se operan solos. Son confirmación dentro de un contexto.",
             font_size=18, color=GRAY_LIGHT)

# Three conditions as cards
add_card(slide, Inches(0.5), Inches(2.3), Inches(3.8), Inches(3.5),
         "ZONA CLAVE", ["Soporte, resistencia", "u otra zona relevante.", "", "Sin zona = sin trade."],
         accent_color=GOLD, title_size=16, body_size=14)

add_card(slide, Inches(4.7), Inches(2.3), Inches(3.8), Inches(3.5),
         "HORA RELEVANTE", ["Esperar el cierre de", "la vela M5.", "", "La hora importa tanto", "como la zona."],
         accent_color=TEAL, title_size=16, body_size=14)

add_card(slide, Inches(8.9), Inches(2.3), Inches(3.8), Inches(3.5),
         "PATRÓN DE VELA", ["Envolvente, martillo,", "rechazo, doji...", "", "Solo valen en zona +", "hora correcta."],
         accent_color=GREEN_CANDLE, title_size=16, body_size=14)

# Bottom emphasis
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
             "ZONA + HORA + PATRÓN = Confluencia de confirmación",
             font_size=16, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

slide.notes_slide.notes_text_frame.text = (
    "Los patrones de velas no se operan solos. Los uso como confirmación cuando están "
    "en una zona importante y en una hora relevante.\n\n"
    "Patrones como los vistos ayer tienen sentido solo si aparecen en zonas clave.\n"
    "La hora también importa: debemos ver cómo cierra la vela de M5."
)


# ============================================================
# SLIDE 9: EL CIERRE DE M5
# ============================================================
slide = next_slide()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "El Cierre de la Vela M5", font_size=32, bold=True, color=WHITE)

# Main quote
quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.4), Inches(11), Inches(1.2))
quote.fill.solid()
quote.fill.fore_color.rgb = CARD_BG
quote.line.color.rgb = GOLD
quote.line.width = Pt(1)

add_text_box(slide, Inches(1.2), Inches(1.55), Inches(10.2), Inches(1.0),
             '"El cierre de la vela M5 confirma la intención del mercado.\n'
             'No entramos mientras la vela se está formando. Esperamos ver cómo cierra."',
             font_size=16, color=GOLD_LIGHT)

# Two scenarios side by side
add_card(slide, Inches(0.5), Inches(3.2), Inches(5.5), Inches(3.0),
         "✓  CIERRE CON FUERZA", ["Cuerpo completo y definido", "Confirma movimiento", "", "→ ENTRADA VÁLIDA"],
         accent_color=GREEN_CANDLE, title_size=16, body_size=15)

add_card(slide, Inches(6.5), Inches(3.2), Inches(5.5), Inches(3.0),
         "✗  CIERRE CON RECHAZO", ["Mecha larga o cambio de color", "Muestra trampa / indecisión", "", "→ NO ENTRAR"],
         accent_color=RED_CANDLE, title_size=16, body_size=15)

slide.notes_slide.notes_text_frame.text = (
    "Siempre hay que esperar ver cómo termina la vela de M5 antes de decidir si entrar o no.\n"
    "A veces parece que el precio va a romper un nivel, pero al cierre la vela cambia "
    "completamente y deja una mecha.\n\n"
    "El cierre es lo que realmente confirma si hay intención o solo fue una trampa.\n\n"
    "• Si cierra con fuerza y cuerpo completo → confirma movimiento.\n"
    "• Si cierra con mecha o cambio de color → muestra rechazo.\n\n"
    "La decisión de entrada se toma SOLO después del cierre."
)


# ============================================================
# SLIDE 10: DIVISOR - PARTE 04
# ============================================================
add_divider_slide("Temporalidades de Análisis", 4)
slide_counter[0] += 1

# ============================================================
# SLIDE 11: MÉTODO H1 → M15 → M1
# ============================================================
slide = next_slide()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Mi Método: H1 → M15 → M1", font_size=32, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Cada temporalidad cumple una función específica en el análisis.",
             font_size=16, color=GRAY_LIGHT)

# Three cards horizontal
add_card(slide, Inches(0.3), Inches(2.0), Inches(4.0), Inches(4.5),
         "H1 · ESTRUCTURA", [
             "Defino la dirección",
             "principal del mercado.",
             "",
             "¿Impulso o retroceso?",
             "¿Alcista o bajista?",
             "",
             "→ TENDENCIA"
         ], accent_color=GOLD, title_size=16, body_size=14)

add_card(slide, Inches(4.6), Inches(2.0), Inches(4.0), Inches(4.5),
         "M15 · CONTEXTO", [
             "¿Cómo se comporta el",
             "precio dentro de H1?",
             "",
             "¿Está alineado con",
             "la estructura mayor?",
             "",
             "→ AJUSTE"
         ], accent_color=TEAL, title_size=16, body_size=14)

add_card(slide, Inches(8.9), Inches(2.0), Inches(4.0), Inches(4.5),
         "M1 · PRECISIÓN", [
             "Marco la línea exacta",
             "donde quiero entrar.",
             "",
             "Solo ejecuto si",
             "TODO coincide.",
             "",
             "→ ENTRADA"
         ], accent_color=GREEN_CANDLE, title_size=16, body_size=14)

slide.notes_slide.notes_text_frame.text = (
    "Explica tu método paso a paso:\n\n"
    "1. H1 – Estructura y tendencia:\n"
    "   'Aquí defino si el mercado está en impulso o retroceso. Identifico la dirección principal.'\n\n"
    "2. M15 – Ajuste y contexto:\n"
    "   'Busco cómo se comporta el precio dentro de esa estructura. ¿Está alineado o en corrección?'\n\n"
    "3. M1 – Precisión de entrada:\n"
    "   'Aquí marco la línea exacta donde quiero entrar. Solo ejecuto si todo coincide.'"
)


# ============================================================
# SLIDE 12: DIVISOR - PARTE 05
# ============================================================
add_divider_slide("Línea de Tendencia", 5)
slide_counter[0] += 1

# ============================================================
# SLIDE 13: QUÉ ES Y CÓMO SE TRAZA
# ============================================================
slide = next_slide()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "La Línea de Tendencia", font_size=32, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             "Marca la dirección dominante del precio. Operamos a su favor.",
             font_size=17, color=GRAY_LIGHT)

# Two cards: alcista vs bajista
add_card(slide, Inches(0.5), Inches(2.2), Inches(5.8), Inches(3.8),
         "TENDENCIA ALCISTA ↗", [
             "Se traza uniendo mínimos",
             "crecientes (pisos más altos).",
             "",
             "La línea va POR DEBAJO",
             "de los puntos de apoyo.",
             "",
             "El precio sube haciendo",
             "retrocesos cada vez más altos."
         ], accent_color=GREEN_CANDLE, title_size=17, body_size=14)

add_card(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(3.8),
         "TENDENCIA BAJISTA ↘", [
             "Se traza uniendo máximos",
             "decrecientes (techos más bajos).",
             "",
             "La línea va POR ENCIMA",
             "de los puntos de rechazo.",
             "",
             "El precio baja haciendo",
             "rebotes cada vez más bajos."
         ], accent_color=RED_CANDLE, title_size=17, body_size=14)

slide.notes_slide.notes_text_frame.text = (
    "Una línea de tendencia se dibuja uniendo mínimos crecientes cuando el mercado "
    "está al alza, o máximos decrecientes cuando está a la baja.\n\n"
    "Sirve para marcar la dirección dominante del precio. Es la que seguimos para "
    "operar a favor de la estructura general.\n\n"
    "'Si el precio cada vez hace pisos más altos, trazamos una línea por debajo "
    "—esa es la línea de tendencia alcista. Si hace techos más bajos, trazamos "
    "una línea por encima —esa es la línea de tendencia bajista.'"
)


# ============================================================
# SLIDE 14: ROMPIMIENTO + RETESTEO
# ============================================================
slide = next_slide()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Rompimiento y Retesteo", font_size=32, bold=True, color=WHITE)

# Main concept box
concept = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.3), Inches(11), Inches(1.5))
concept.fill.solid()
concept.fill.fore_color.rgb = CARD_BG
concept.line.color.rgb = GOLD
concept.line.width = Pt(1)

add_text_box(slide, Inches(1.2), Inches(1.5), Inches(10.2), Inches(1.2),
             "Cuando el precio rompe su línea de tendencia, NO significa\n"
             "automáticamente que cambió la dirección del mercado.",
             font_size=17, bold=True, color=WHITE)

# Steps
add_multi_text(slide, Inches(0.8), Inches(3.2), Inches(11), Inches(3.5),
    ["1.  El precio ROMPE la línea de tendencia",
     "",
     "2.  Crea una contratendencia temporal",
     "",
     "3.  VUELVE A TESTEAR la línea rota (desde el otro lado)",
     "",
     "4.  Si RESPETA esa línea rota → confirma continuación o cambio estructural"],
    font_size=17, color=WHITE)

# Emphasis at bottom
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
             "El retesteo desde el lado contrario = momento clave de decisión",
             font_size=15, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)

slide.notes_slide.notes_text_frame.text = (
    "Cuando el precio rompe su línea de tendencia, no significa automáticamente "
    "que cambió la dirección del mercado; muchas veces, lo que está haciendo es "
    "crear una contratendencia para luego volver a testear la línea rota desde el otro lado.\n\n"
    "Ese testeo es clave: si el precio respeta esa línea rota, ahora desde el lado contrario, "
    "suele marcar el momento donde la tendencia principal puede continuar o confirmarse "
    "el cambio estructural."
)


# ============================================================
# SLIDE 15: RESUMEN DEL DÍA
# ============================================================
slide = next_slide()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Resumen de la Clase", font_size=32, bold=True, color=WHITE)

# Summary points as mini cards
topics = [
    ("01", "LECTURA DEL PRECIO", "HH/HL = alcista · LH/LL = bajista", GOLD),
    ("02", "SOPORTES Y RESISTENCIAS", "3 testeos: descubre → valida → operativo", TEAL),
    ("03", "CONFIRMACIÓN CON VELAS", "Zona + Hora + Patrón = Confluencia", GREEN_CANDLE),
    ("04", "TEMPORALIDADES", "H1 estructura → M15 contexto → M1 entrada", GOLD),
    ("05", "LÍNEA DE TENDENCIA", "Trazar a favor + Rompimiento y retesteo", TEAL),
]

y_pos = Inches(1.4)
for num, title, desc, color in topics:
    # Number
    add_text_box(slide, Inches(0.8), y_pos, Inches(0.8), Inches(0.8),
                 num, font_size=22, bold=True, color=color)
    # Title
    add_text_box(slide, Inches(1.8), y_pos, Inches(4), Inches(0.5),
                 title, font_size=15, bold=True, color=WHITE)
    # Description
    add_text_box(slide, Inches(1.8), y_pos + Inches(0.4), Inches(9), Inches(0.4),
                 desc, font_size=13, color=GRAY_LIGHT)
    y_pos += Inches(1.05)

slide.notes_slide.notes_text_frame.text = (
    "Resumen:\n"
    "Hoy aprendimos a leer el precio (estructura HH/HL), identificar zonas clave "
    "(soportes/resistencias con método de testeos), confirmar con velas y el cierre M5, "
    "usar las temporalidades H1→M15→M1, y trazar líneas de tendencia con su rompimiento y retesteo."
)

# ============================================================
# SLIDE 16: TAREA
# ============================================================
slide = next_slide()

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "Tarea", font_size=32, bold=True, color=GOLD)

add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5),
    ["1.  Abrir un gráfico (V75 o EUR/USD) y marcar:",
     "     • Soportes y resistencias con sus testeos",
     "     • Estructura (HH/HL o LH/LL)",
     "     • Línea de tendencia principal",
     "",
     "2.  Buscar 2 ejemplos donde el cierre M5 confirmó",
     "     movimiento y 2 donde mostró rechazo (trampa).",
     "",
     "3.  Identificar en qué temporalidad (H1/M15/M1)",
     "     se ve mejor cada elemento marcado.",
     "",
     "4.  Repasar lo visto hoy."],
    font_size=17, color=WHITE)

slide.notes_slide.notes_text_frame.text = (
    "Tarea para mañana:\n"
    "1. Abrir un gráfico y marcar S/R, estructura y línea de tendencia.\n"
    "2. Buscar ejemplos de cierre M5 que confirma vs. que rechaza.\n"
    "3. Identificar temporalidades.\n"
    "4. Repasar lo visto hoy."
)

# ============================================================
# SAVE
# ============================================================
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "Clase_03_PREMIUM.pptx")
prs.save(output_path)
print(f"✓ Generado: {output_path}")
print(f"  Total diapositivas: {len(prs.slides)}")
